from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand palette ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # slide background
ELECTRIC   = RGBColor(0x00, 0xC2, 0xFF)   # accent / hero numbers
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
ORANGE     = RGBColor(0xFF, 0x6B, 0x35)   # warning / pain-point red-flag
GREEN      = RGBColor(0x00, 0xE6, 0x96)   # positive outcome

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helpers ──────────────────────────────────────────────────────────────────
def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, text, font_size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, bg_color=None, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    return txBox

def rule(slide, t, color=ELECTRIC, l=0.4, w=12.5, h=0.04):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def card(slide, l, t, w, h, fill=RGBColor(0x12, 0x28, 0x3A)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def bullet_box(slide, l, t, w, h, lines, size=16, color=WHITE, spacing=0.42):
    """Render a list of strings as bullet rows."""
    for i, line in enumerate(lines):
        box(slide, l, t + i * spacing, w, spacing + 0.1, line,
            font_size=size, color=color)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
# Accent stripe left edge
stripe = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
stripe.fill.solid(); stripe.fill.fore_color.rgb = ELECTRIC; stripe.line.fill.background()

box(s, 0.55, 1.0, 12.0, 0.7, "SOLVING YOUR CASH APPLICATION PROBLEM",
    font_size=13, color=ELECTRIC, bold=True)
box(s, 0.55, 1.65, 11.5, 1.6,
    "Automating Payment Matching Across\nEvery Incoming Channel",
    font_size=44, bold=True, color=WHITE)
rule(s, 3.6)
box(s, 0.55, 3.75, 6.0, 0.5, "Prepared for:  Peter Sukits, Senior Finance Leader",
    font_size=16, color=LIGHT_GRAY)
box(s, 0.55, 4.25, 6.0, 0.5, "LunarLogic  ·  Jonathan Rodriguez  ·  May 2026",
    font_size=14, color=LIGHT_GRAY)

# Right side hero stat
card(s, 8.4, 2.8, 4.5, 3.5)
box(s, 8.5, 3.0, 4.3, 0.55, "Clients using LunarLogic see",
    font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
box(s, 8.5, 3.45, 4.3, 1.1, "19-Day",
    font_size=56, bold=True, color=ELECTRIC, align=PP_ALIGN.CENTER)
box(s, 8.5, 4.45, 4.3, 0.55, "average DSO reduction",
    font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 8.5, 5.05, 4.3, 0.4, "— Kaptain Clean LLC, anchor client",
    font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PAIN POINT (empathy slide)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.7, "THE PROBLEM WE HEARD",
    font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.8,
    "Payments Arrive From Everywhere — Matching Them Is a Manual Nightmare",
    font_size=28, bold=True, color=WHITE)
rule(s, 1.8)

pain_points = [
    ("ACH / Wire transfers", "Reference fields are truncated or blank — impossible to auto-match."),
    ("Check payments",       "Remittance detail lives in a scanned PDF, not your ERP."),
    ("Credit card / portal", "Customer portals post lump-sum payments against multiple invoices."),
    ("Partial payments",     "Split payments across invoices require manual allocation decisions."),
    ("High volume",          "Finance teams spend hours each week on cash application instead of analysis."),
]

for i, (channel, desc) in enumerate(pain_points):
    top = 2.05 + i * 0.92
    card(s, 0.4, top, 5.5, 0.78)
    card(s, 6.1, top, 6.8, 0.78, fill=RGBColor(0x0F, 0x22, 0x33))
    box(s, 0.55, top + 0.08, 5.2, 0.6, f"⚠  {channel}",
        font_size=15, bold=True, color=ORANGE)
    box(s, 6.2, top + 0.08, 6.5, 0.6, desc, font_size=14, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — COST OF THE STATUS QUO
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.6, "WHY THIS MATTERS", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.8,
    "Every Day of Delay in Cash Application Costs Real Money",
    font_size=28, bold=True, color=WHITE)
rule(s, 1.75)

stats = [
    ("4–8 hrs",  "per week",    "staff time lost\nto manual matching"),
    ("3–5 days", "average lag", "between payment receipt\nand QB application"),
    ("$12–40K",  "per year",    "fully-loaded cost of\ncash application errors & rework"),
    ("DSO +5",   "hidden days", "inflated by unapplied\ncash sitting on account"),
]
for i, (num, label, desc) in enumerate(stats):
    l = 0.35 + i * 3.18
    card(s, l, 2.1, 3.0, 4.5)
    box(s, l + 0.12, 2.3, 2.8, 1.0, num, font_size=42, bold=True,
        color=ORANGE, align=PP_ALIGN.CENTER)
    box(s, l + 0.12, 3.15, 2.8, 0.45, label, font_size=13,
        color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    rule(s, 3.7, color=ELECTRIC, l=l+0.3, w=2.4, h=0.03)
    box(s, l + 0.1, 3.85, 2.85, 0.9, desc, font_size=13,
        color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.4, 6.7, 12.5, 0.4,
    "Note: ranges based on professional services firms with 8–20 employees and 100–400 invoices/month.",
    font_size=10, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE LUNARLOGIC SOLUTION OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.6, "OUR SOLUTION", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "LunarLogic Payment Receipt & Cash Application — Automated End-to-End",
    font_size=27, bold=True, color=WHITE)
rule(s, 1.78)

# Flow steps
steps = [
    ("1", "CAPTURE",    "Plaid webhook\ncatches every\nbank transaction\nin real time"),
    ("2", "PARSE",      "AI reads\nremittance data\nfrom any format\n(ACH, check, portal)"),
    ("3", "MATCH",      "Fuzzy-match engine\ncorrelates payment\nto open QuickBooks\ninvoices"),
    ("4", "AUTO-APPLY", "Confidence ≥ 90%:\napplied instantly\nto QB with zero\nhuman touch"),
    ("5", "ESCALATE",   "Ambiguous matches:\nSlack prompt to\nfinance team for\none-click approval"),
    ("6", "LOG",        "Every action logged\nto Google Sheets\nfor audit trail &\nreconciliation"),
]
for i, (num, title, desc) in enumerate(steps):
    l = 0.3 + i * 2.16
    card(s, l, 2.0, 2.0, 4.6)
    box(s, l+0.08, 2.1, 1.85, 0.65, num, font_size=32, bold=True,
        color=ELECTRIC, align=PP_ALIGN.CENTER)
    box(s, l+0.08, 2.65, 1.85, 0.55, title, font_size=13, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    rule(s, 3.3, color=ELECTRIC, l=l+0.2, w=1.65, h=0.03)
    box(s, l+0.05, 3.42, 1.92, 1.9, desc, font_size=12,
        color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Arrow connectors (simple chevrons as text)
for i in range(5):
    l = 2.22 + i * 2.16
    box(s, l, 3.85, 0.22, 0.4, "▶", font_size=14, color=ELECTRIC, align=PP_ALIGN.CENTER)

box(s, 0.4, 6.7, 12.5, 0.4,
    "Default rule: payments applied to oldest open invoice first.  Fully configurable.",
    font_size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW THE MATCHING ENGINE WORKS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.6, "UNDER THE HOOD", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Intelligent Matching — Handles the Edge Cases That Break Manual Processes",
    font_size=27, bold=True, color=WHITE)
rule(s, 1.78)

# Left: matching logic
card(s, 0.35, 1.95, 6.0, 5.15)
box(s, 0.55, 2.1, 5.6, 0.5, "MATCHING SIGNALS USED", font_size=12,
    bold=True, color=ELECTRIC)
signals = [
    "✓  Invoice number extracted from remittance memo",
    "✓  Dollar amount ± partial payment tolerance",
    "✓  Customer name fuzzy match (handles abbreviations)",
    "✓  Payment date vs. invoice due date proximity",
    "✓  PO number or job reference cross-reference",
    "✓  Historical payment behavior per customer",
]
for i, sig in enumerate(signals):
    box(s, 0.55, 2.7 + i * 0.55, 5.6, 0.5, sig, font_size=14, color=WHITE)

# Right: confidence tiers
card(s, 6.6, 1.95, 6.35, 5.15)
box(s, 6.8, 2.1, 6.0, 0.5, "CONFIDENCE TIERS & ACTIONS", font_size=12,
    bold=True, color=ELECTRIC)

tiers = [
    (GREEN,      "≥ 90% confidence",    "Auto-apply to QuickBooks immediately"),
    (ELECTRIC,   "70–89% confidence",   "Apply + flag for next-day review"),
    (ORANGE,     "< 70% confidence",    "Slack prompt — one-click approval"),
    (RGBColor(0x99,0x00,0x00), "Unmatched", "Dead-letter queue + Slack alert"),
]
for i, (color, tier, action) in enumerate(tiers):
    top = 2.75 + i * 1.05
    dot = s.shapes.add_shape(9, Inches(6.8), Inches(top), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
    box(s, 7.15, top - 0.05, 5.5, 0.38, tier, font_size=14, bold=True, color=color)
    box(s, 7.15, top + 0.33, 5.5, 0.38, action, font_size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PROOF POINT
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.6, "PROOF POINT", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "What LunarLogic Automation Delivers — A Client Story",
    font_size=28, bold=True, color=WHITE)
rule(s, 1.78)

# Quote block
card(s, 0.35, 1.95, 12.6, 2.1, fill=RGBColor(0x08, 0x1A, 0x28))
box(s, 0.65, 2.05, 12.1, 0.45, "KAPTAIN CLEAN LLC  ·  Commercial Cleaning Services",
    font_size=11, bold=True, color=ELECTRIC)
box(s, 0.55, 2.45, 12.1, 1.3,
    '"Before LunarLogic, our AP team spent half a day every week chasing down\n'
    'which checks applied to which invoice. Now it just happens. Our books close faster\n'
    'and we stopped carrying receivables we didn\'t know were already paid."',
    font_size=14, color=WHITE)

# Result cards
results = [
    ("84%",        "reduction in\ninvoice processing time"),
    ("19 Days",    "DSO improvement\npost go-live"),
    ("< 2 Weeks",  "time to full\nproduction deployment"),
    ("0 hrs/wk",   "manual cash\napplication time"),
]
for i, (num, label) in enumerate(results):
    l = 0.35 + i * 3.18
    card(s, l, 4.3, 3.0, 2.8)
    box(s, l+0.1, 4.5, 2.8, 1.0, num, font_size=40, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)
    box(s, l+0.1, 5.35, 2.8, 0.9, label, font_size=14,
        color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT YOU GET (integrations & deliverables)
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.6, "WHAT YOU GET", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Everything Included — No Bolt-Ons, No Hidden Fees",
    font_size=28, bold=True, color=WHITE)
rule(s, 1.78)

deliverables = [
    ("Plaid Bank Integration",        "Real-time transaction feed from your bank accounts"),
    ("QuickBooks Online Sync",         "Payments auto-posted and matched to open invoices"),
    ("AI Matching Engine",             "Trained on your customer and invoice history"),
    ("Slack Approval Workflow",        "One-click resolution for ambiguous matches"),
    ("AR Aging Dashboard",             "Live React dashboard — DSO trend, aging buckets, invoice status"),
    ("Full Audit Log",                 "Every match decision logged to Google Sheets"),
    ("Daily AR Summary to Slack",      "Finance team sees AR health every morning at 9 AM"),
    ("Dedicated Implementation Lead",  "White-glove setup — live in under 2 weeks"),
]
for i, (title, desc) in enumerate(deliverables):
    row = i % 4
    col = i // 4
    l = 0.35 + col * 6.55
    t = 2.05 + row * 1.2
    card(s, l, t, 6.2, 1.05)
    box(s, l+0.2, t+0.08, 5.8, 0.42, f"◆  {title}", font_size=14, bold=True, color=ELECTRIC)
    box(s, l+0.2, t+0.52, 5.8, 0.42, desc, font_size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — IMPLEMENTATION TIMELINE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.6, "IMPLEMENTATION", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Live in Under Two Weeks — Zero Disruption to Your Current Process",
    font_size=28, bold=True, color=WHITE)
rule(s, 1.78)

phases = [
    ("Days 1–2",   "KICKOFF & CONNECT",
     ["QuickBooks OAuth credentialing",
      "Plaid bank account linking",
      "Slack workspace integration",
      "Map your chart of accounts"]),
    ("Days 3–5",   "CONFIGURE & TRAIN",
     ["Load customer master & invoice history",
      "Configure confidence thresholds",
      "Set default allocation rules",
      "Define escalation contacts"]),
    ("Days 6–10",  "TEST & VALIDATE",
     ["Parallel-run against your last 30 days",
      "Review match accuracy report",
      "Tune edge cases with your team",
      "Finance sign-off on match logic"]),
    ("Day 11–14",  "GO LIVE",
     ["Flip switch to production",
      "Decommission manual process",
      "Live monitoring for first week",
      "Dashboard access for your team"]),
]
for i, (period, title, items) in enumerate(phases):
    l = 0.35 + i * 3.22
    card(s, l, 1.95, 3.05, 5.15)
    box(s, l+0.12, 2.05, 2.82, 0.42, period, font_size=12, color=ELECTRIC, bold=True)
    box(s, l+0.12, 2.48, 2.82, 0.52, title, font_size=14, bold=True, color=WHITE)
    rule(s, 3.1, color=ELECTRIC, l=l+0.2, w=2.65, h=0.03)
    for j, item in enumerate(items):
        box(s, l+0.12, 3.22 + j * 0.67, 2.82, 0.6,
            f"→  {item}", font_size=12, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — INVESTMENT / PRICING
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3, 12.5, 0.6, "INVESTMENT", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Transparent, Flat-Rate Pricing — No Per-Transaction Fees",
    font_size=28, bold=True, color=WHITE)
rule(s, 1.78)

plans = [
    ("Essentials",    "$697",  "/mo", "Up to 150 invoices/month",
     ["Full cash application automation",
      "QuickBooks + Plaid integration",
      "Slack approval workflow",
      "AR aging dashboard",
      "Email support"]),
    ("Professional",  "$1,497", "/mo", "Up to 250 invoices/month",
     ["Everything in Essentials",
      "Priority Slack support",
      "Custom matching rules",
      "Advanced reporting",
      "Dedicated CSM"]),
    ("Business",      "$2,497", "/mo", "Up to 400 invoices/month",
     ["Everything in Professional",
      "Multi-entity support",
      "API access",
      "Custom integrations",
      "SLA guarantee"]),
]

highlight = [False, True, False]
for i, (name, price, per, vol, feats) in enumerate(plans):
    l = 0.5 + i * 4.2
    fill = RGBColor(0x00, 0x3A, 0x5C) if highlight[i] else RGBColor(0x12, 0x28, 0x3A)
    card(s, l, 1.92, 3.9, 5.25, fill=fill)
    if highlight[i]:
        box(s, l+0.1, 1.92, 3.7, 0.38, "  MOST POPULAR",
            font_size=11, bold=True, color=NAVY, bg_color=ELECTRIC, align=PP_ALIGN.CENTER)
    box(s, l+0.15, 2.45, 3.6, 0.5, name, font_size=18, bold=True, color=WHITE)
    box(s, l+0.15, 2.92, 2.4, 0.8, price, font_size=38, bold=True, color=ELECTRIC)
    box(s, l+2.3, 3.35, 1.2, 0.42, per, font_size=14, color=LIGHT_GRAY)
    box(s, l+0.15, 3.72, 3.6, 0.38, vol, font_size=12, color=LIGHT_GRAY)
    rule(s, 4.2, color=ELECTRIC, l=l+0.3, w=3.25, h=0.03)
    for j, feat in enumerate(feats):
        box(s, l+0.15, 4.35 + j * 0.47, 3.65, 0.42,
            f"✓  {feat}", font_size=12, color=WHITE)

box(s, 0.4, 7.0, 12.5, 0.38,
    "Implementation fee: $2,500  (waived with 12-month commitment)  ·  Overage: $5/invoice  ·  60-day satisfaction guarantee",
    font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — NEXT STEPS / CTA
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
stripe2 = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
stripe2.fill.solid(); stripe2.fill.fore_color.rgb = ELECTRIC; stripe2.line.fill.background()

box(s, 0.55, 0.9, 12.0, 0.6, "NEXT STEPS", font_size=11, color=ELECTRIC, bold=True)
box(s, 0.55, 1.4, 11.5, 1.0,
    "Let's Get Your Cash Application Running on Autopilot",
    font_size=32, bold=True, color=WHITE)
rule(s, 2.6)

steps_next = [
    ("This Week",       "30-minute technical scoping call — map your bank accounts, QB setup, and invoice volume"),
    ("Within 48 Hours", "Custom proposal with recommended plan tier and ROI projection for your firm"),
    ("Day 1 of Onboarding", "Credentials exchange and Plaid/QB connection — we handle all the technical setup"),
    ("Day 14",          "Go live — your team wakes up to cash application that runs itself"),
]
for i, (when, what) in enumerate(steps_next):
    t = 2.85 + i * 0.98
    dot = s.shapes.add_shape(9, Inches(0.55), Inches(t+0.1), Inches(0.25), Inches(0.25))
    dot.fill.solid(); dot.fill.fore_color.rgb = ELECTRIC; dot.line.fill.background()
    box(s, 0.9, t, 2.8, 0.45, when, font_size=14, bold=True, color=ELECTRIC)
    box(s, 0.9, t + 0.43, 11.5, 0.45, what, font_size=14, color=LIGHT_GRAY)

# Contact card
card(s, 7.5, 2.75, 5.45, 4.0)
box(s, 7.7, 2.95, 5.1, 0.5, "JONATHAN RODRIGUEZ", font_size=16, bold=True, color=WHITE)
box(s, 7.7, 3.45, 5.1, 0.45, "Founder & CEO, LunarLogic", font_size=13, color=LIGHT_GRAY)
rule(s, 4.05, color=ELECTRIC, l=7.7, w=4.8, h=0.03)
box(s, 7.7, 4.2, 5.1, 0.45, "✉  jrodriguez@lunarlogic.ai", font_size=14, color=WHITE)
box(s, 7.7, 4.75, 5.1, 0.45, "60-day satisfaction guarantee", font_size=13, color=GREEN)
box(s, 7.7, 5.25, 5.1, 0.8,
    "If we don't reduce your cash application\ntime by 80%, you pay nothing.",
    font_size=13, color=LIGHT_GRAY)

# ── Save ─────────────────────────────────────────────────────────────────
out = "/home/user/lunarlogic-dashboard/LunarLogic_CashApplication_PeterSukits.pptx"
prs.save(out)
print(f"Saved → {out}")
