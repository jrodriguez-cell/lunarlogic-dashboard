from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
ELECTRIC   = RGBColor(0x00, 0xC2, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
ORANGE     = RGBColor(0xFF, 0x6B, 0x35)
GREEN      = RGBColor(0x00, 0xE6, 0x96)
DARK_CARD  = RGBColor(0x12, 0x28, 0x3A)
DARKER     = RGBColor(0x08, 0x1A, 0x28)
MID_BLUE   = RGBColor(0x00, 0x3A, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, text, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, bg_color=None, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    if bg_color:
        tb.fill.solid(); tb.fill.fore_color.rgb = bg_color
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb

def rule(slide, t, color=ELECTRIC, l=0.4, w=12.5, h=0.04):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def card(slide, l, t, w, h, fill=DARK_CARD):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

def dot(slide, l, t, size=0.22, color=ELECTRIC):
    sh = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(size), Inches(size))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def left_stripe(slide):
    sh = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    sh.fill.solid(); sh.fill.fore_color.rgb = ELECTRIC; sh.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); left_stripe(s)

box(s, 0.55, 1.0,  12.0, 0.6,  "CASH APPLICATION AUTOMATION",
    size=12, color=ELECTRIC, bold=True)
box(s, 0.55, 1.55, 11.5, 1.8,
    "Eliminating the Manual Work\nBetween Payment Receipt\nand Closed Receivables",
    size=38, bold=True, color=WHITE)
rule(s, 3.75)
box(s, 0.55, 3.92, 7.0, 0.5,
    "Prepared for:  Peter Sukits, Senior Finance Leader",
    size=16, color=LIGHT_GRAY)
box(s, 0.55, 4.42, 7.0, 0.5,
    "LunarLogic  ·  Jonathan Rodriguez  ·  May 2026",
    size=14, color=LIGHT_GRAY)

card(s, 8.5, 2.9, 4.4, 3.3)
box(s, 8.6, 3.1,  4.2, 0.5,  "Clients see an average of",
    size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
box(s, 8.6, 3.52, 4.2, 1.05, "19-Day",
    size=54, bold=True, color=ELECTRIC, align=PP_ALIGN.CENTER)
box(s, 8.6, 4.5,  4.2, 0.5,  "DSO reduction after go-live",
    size=16, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 8.6, 5.05, 4.2, 0.4,  "— Kaptain Clean LLC, anchor client",
    size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "THE PROBLEM", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Payments Arrive From Everywhere — Matching Them Is a Manual Nightmare",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

pains = [
    ("ACH & Wire Transfers",
     "Reference fields are truncated or missing — impossible to auto-match without manual investigation."),
    ("Check Payments",
     "Remittance detail lives on a paper stub or scanned PDF, completely disconnected from your ERP."),
    ("Customer Portal & Card Payments",
     "Portals post lump-sum settlements against multiple invoices — allocations must be untangled manually."),
    ("Partial & Split Payments",
     "Short-pays and installments require judgment calls on which invoice to apply against first."),
    ("Volume & Velocity",
     "Finance teams lose 4–8 hours per week on cash application — time that should go to analysis, not data entry."),
]
for i, (title, desc) in enumerate(pains):
    t = 2.05 + i * 0.92
    card(s, 0.4,  t, 5.5, 0.78)
    card(s, 6.1,  t, 6.8, 0.78, fill=RGBColor(0x0F, 0x22, 0x33))
    box(s, 0.55, t+0.08, 5.2, 0.6, f"⚠  {title}",  size=14, bold=True, color=ORANGE)
    box(s, 6.2,  t+0.08, 6.5, 0.62, desc, size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — COST OF THE STATUS QUO
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "THE REAL COST", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Manual Cash Application Is a Revenue and Efficiency Problem",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

stats = [
    ("4–8 hrs",  "per week",       "staff time consumed\nby manual matching"),
    ("3–5 days", "average lag",    "from payment receipt\nto ledger close"),
    ("$15–40K",  "annual cost",    "fully-loaded cost of\nerrors, rework & delays"),
    ("+5 Days",  "hidden DSO",     "from unapplied cash\nsitting in suspense"),
]
for i, (num, label, desc) in enumerate(stats):
    l = 0.35 + i * 3.18
    card(s, l, 2.1, 3.0, 4.55)
    box(s, l+0.12, 2.3,  2.8, 1.0,  num,   size=42, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    box(s, l+0.12, 3.18, 2.8, 0.45, label, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    rule(s, 3.73, color=ELECTRIC, l=l+0.3, w=2.4, h=0.03)
    box(s, l+0.1,  3.88, 2.85, 0.9, desc,  size=13, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.4, 6.9, 12.5, 0.38,
    "Benchmarks based on professional services firms with 8–20 employees and 100–400 invoices/month.",
    size=10, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION OVERVIEW (flow)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "THE SOLUTION", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "LunarLogic Cash Application — Automated End-to-End",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

steps = [
    ("1", "CAPTURE",     "Real-time bank\nfeed pulls every\nincoming payment\nautomatically"),
    ("2", "PARSE",       "AI extracts\nremittance data\nfrom any format\nor channel"),
    ("3", "MATCH",       "Fuzzy-match\nengine links\npayments to\nopen invoices"),
    ("4", "AUTO-APPLY",  "High-confidence\nmatches post to\nyour ERP with\nzero human touch"),
    ("5", "ESCALATE",    "Ambiguous items\ntrigger a one-click\napproval prompt\nto your team"),
    ("6", "LOG",         "Every decision\nlogged for full\naudit trail &\nreconciliation"),
]
for i, (num, title, desc) in enumerate(steps):
    l = 0.3 + i * 2.16
    card(s, l, 2.0, 2.0, 4.65)
    box(s, l+0.08, 2.1,  1.85, 0.65, num,   size=32, bold=True, color=ELECTRIC, align=PP_ALIGN.CENTER)
    box(s, l+0.08, 2.68, 1.85, 0.55, title, size=12, bold=True, color=WHITE,    align=PP_ALIGN.CENTER)
    rule(s, 3.33, color=ELECTRIC, l=l+0.2, w=1.65, h=0.03)
    box(s, l+0.05, 3.44, 1.92, 1.95, desc,  size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

for i in range(5):
    box(s, 2.23 + i * 2.16, 3.88, 0.22, 0.4, "▶", size=14, color=ELECTRIC, align=PP_ALIGN.CENTER)

box(s, 0.4, 6.88, 12.5, 0.4,
    "Default rule: apply payment to oldest open invoice first.  All rules are configurable.",
    size=11, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — MATCHING ENGINE
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "HOW MATCHING WORKS", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Intelligent Matching That Handles the Edge Cases Breaking Your Team Today",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

card(s, 0.35, 1.95, 6.0, 5.2)
box(s, 0.55, 2.1, 5.6, 0.5, "SIGNALS THE ENGINE USES", size=12, bold=True, color=ELECTRIC)
signals = [
    "✓  Invoice or PO number extracted from remittance memo",
    "✓  Payment amount vs. invoice total (partial tolerance built in)",
    "✓  Payer name — fuzzy match handles abbreviations & variations",
    "✓  Payment date relative to invoice due date",
    "✓  Job reference, contract number, or custom identifiers",
    "✓  Historical payment behavior learned per customer over time",
]
for i, sig in enumerate(signals):
    box(s, 0.55, 2.72 + i * 0.55, 5.6, 0.5, sig, size=14, color=WHITE)

card(s, 6.6, 1.95, 6.35, 5.2)
box(s, 6.8, 2.1, 6.0, 0.5, "CONFIDENCE TIERS & ACTIONS", size=12, bold=True, color=ELECTRIC)

tiers = [
    (GREEN,                        "≥ 90% confidence",   "Post to ERP automatically — no human required"),
    (ELECTRIC,                     "70–89% confidence",  "Post and flag for next-day finance review"),
    (ORANGE,                       "< 70% confidence",   "Hold and send one-click approval to your team"),
    (RGBColor(0xCC, 0x22, 0x22),   "Unmatched",          "Dead-letter queue — alert sent immediately"),
]
for i, (color, tier, action) in enumerate(tiers):
    t = 2.75 + i * 1.05
    dot(s, 6.8, t+0.08, color=color)
    box(s, 7.18, t,      5.5, 0.42, tier,   size=14, bold=True, color=color)
    box(s, 7.18, t+0.43, 5.5, 0.42, action, size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — ROI CALCULATOR
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "RETURN ON INVESTMENT", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "The Math Is Simple — And It Compounds Every Month",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

# Left — cost inputs
card(s, 0.35, 1.95, 5.9, 5.15)
box(s, 0.55, 2.08, 5.5, 0.5, "TYPICAL COST INPUTS (annual)", size=12, bold=True, color=ELECTRIC)
inputs = [
    ("Staff time on cash application",  "6 hrs/wk  ×  $40/hr  ×  50 wks",  "$12,000"),
    ("Error correction & rework",       "Avg 1.5 hrs/wk  ×  $40/hr",        "$3,000"),
    ("Late payment financing cost",     "5-day lag  ×  $500K AR  ×  6% APR", "$4,100"),
    ("Auditor / reconciliation time",   "10 hrs/mo  ×  $75/hr",              "$9,000"),
    ("TOTAL MANUAL COST",               "",                                   "$28,100+"),
]
for i, (label, detail, value) in enumerate(inputs):
    t = 2.72 + i * 0.82
    is_total = i == 4
    clr = ELECTRIC if is_total else WHITE
    sz  = 14 if is_total else 13
    box(s, 0.55, t, 3.2, 0.38, label, size=sz, bold=is_total, color=clr)
    box(s, 0.55, t+0.38, 3.5, 0.35, detail, size=11, color=LIGHT_GRAY)
    box(s, 4.3, t+0.05, 1.7, 0.45, value, size=sz, bold=is_total, color=clr, align=PP_ALIGN.RIGHT)

# Right — savings
card(s, 6.55, 1.95, 6.4, 5.15)
box(s, 6.75, 2.08, 6.0, 0.5, "LUNARLOGIC OUTCOME", size=12, bold=True, color=ELECTRIC)

outcomes = [
    (GREEN,  "~$25K",  "annual staff time recaptured"),
    (GREEN,  "19 days","average DSO improvement"),
    (GREEN,  "80%+",   "reduction in match exceptions"),
    (GREEN,  "< 24 hrs","average payment-to-close cycle"),
]
for i, (color, num, label) in enumerate(outcomes):
    t = 2.72 + i * 0.95
    card(s, 6.75, t, 6.0, 0.8, fill=RGBColor(0x08, 0x24, 0x18))
    box(s, 6.95, t+0.08, 1.8, 0.62, num,   size=24, bold=True, color=color)
    box(s, 8.7,  t+0.2,  3.8, 0.45, label, size=13, color=WHITE)

# Payback callout
card(s, 6.75, 6.55-0.55, 6.0, 1.05, fill=MID_BLUE)
box(s, 6.95, 6.08-0.55, 5.7, 0.88,
    "At $1,497/mo (Professional), annual investment = $17,964\nTypical client ROI:  56% in year one — before DSO gains.",
    size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — PROOF POINT
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "CLIENT PROOF POINT", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "What Changes When Cash Application Runs on Autopilot",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

card(s, 0.35, 1.95, 12.6, 2.15, fill=DARKER)
box(s, 0.65, 2.05, 12.0, 0.45,
    "KAPTAIN CLEAN LLC  ·  Commercial Services  ·  LunarLogic Client Since 2025",
    size=11, bold=True, color=ELECTRIC)
box(s, 0.55, 2.48, 12.1, 1.35,
    '"Before LunarLogic, our finance team spent half a day every week tracking down which payments\n'
    'applied to which invoice. Now it just happens. Our books close faster and we stopped carrying\n'
    'receivables we didn\'t even know were already paid."',
    size=14, color=WHITE)

results = [
    ("84%",       "reduction in\ninvoice processing time"),
    ("19 Days",   "DSO improvement\npost go-live"),
    ("< 2 Weeks", "from contract\nto production"),
    ("0 hrs/wk",  "manual cash\napplication time"),
]
for i, (num, label) in enumerate(results):
    l = 0.35 + i * 3.18
    card(s, l, 4.35, 3.0, 2.8)
    box(s, l+0.1, 4.55, 2.8, 1.0, num,   size=40, bold=True, color=GREEN,  align=PP_ALIGN.CENTER)
    box(s, l+0.1, 5.42, 2.8, 0.85, label, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — WHAT YOU GET
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "WHAT YOU GET", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "A Complete Solution — No Bolt-Ons, No Hidden Fees",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

deliverables = [
    ("Bank Feed Integration",          "Real-time connection to your financial institution — any bank"),
    ("ERP / Accounting System Sync",   "Payments auto-posted and matched to open invoices in your system"),
    ("AI Matching Engine",             "Trained on your customer history and invoice patterns"),
    ("Team Approval Workflow",         "One-click resolution for exceptions — works in Slack, email, or browser"),
    ("AR Aging Dashboard",             "Live dashboard: DSO trend, aging buckets, invoice status at a glance"),
    ("Full Audit Log",                 "Every match decision recorded — searchable, exportable, audit-ready"),
    ("Daily AR Health Summary",        "Finance team receives AR snapshot each morning — no manual pull"),
    ("Dedicated Implementation Lead",  "White-glove onboarding — live in under two weeks"),
]
for i, (title, desc) in enumerate(deliverables):
    row = i % 4
    col = i // 4
    l = 0.35 + col * 6.55
    t = 2.05 + row * 1.2
    card(s, l, t, 6.2, 1.05)
    box(s, l+0.2, t+0.08, 5.8, 0.42, f"◆  {title}", size=14, bold=True, color=ELECTRIC)
    box(s, l+0.2, t+0.52, 5.8, 0.42, desc, size=13, color=LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — PRICING
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
box(s, 0.4, 0.3,  12.5, 0.6,  "INVESTMENT", size=11, color=ELECTRIC, bold=True)
box(s, 0.4, 0.75, 12.0, 0.85,
    "Transparent Flat-Rate Pricing — No Per-Transaction Fees",
    size=27, bold=True, color=WHITE)
rule(s, 1.78)

plans = [
    ("Essentials",   "$697",   "Up to 150 invoices/month",
     ["Full cash application automation",
      "Bank feed + ERP integration",
      "Team approval workflow",
      "AR aging dashboard",
      "Email support"]),
    ("Professional", "$1,497", "Up to 250 invoices/month",
     ["Everything in Essentials",
      "Priority support channel",
      "Custom matching rules",
      "Advanced reporting",
      "Dedicated customer success"]),
    ("Business",     "$2,497", "Up to 400 invoices/month",
     ["Everything in Professional",
      "Multi-entity support",
      "API access",
      "Custom integrations",
      "SLA guarantee"]),
]
highlight = [False, True, False]
for i, (name, price, vol, feats) in enumerate(plans):
    l = 0.5 + i * 4.2
    fill = MID_BLUE if highlight[i] else DARK_CARD
    card(s, l, 1.92, 3.9, 5.3, fill=fill)
    if highlight[i]:
        box(s, l+0.1, 1.92, 3.7, 0.38, "  MOST POPULAR",
            size=11, bold=True, color=NAVY, bg_color=ELECTRIC, align=PP_ALIGN.CENTER)
    box(s, l+0.15, 2.45, 3.6, 0.52, name,  size=18, bold=True, color=WHITE)
    box(s, l+0.15, 2.95, 2.4, 0.82, price, size=38, bold=True, color=ELECTRIC)
    box(s, l+2.35, 3.38, 1.1, 0.38, "/mo",  size=14, color=LIGHT_GRAY)
    box(s, l+0.15, 3.78, 3.6, 0.38, vol,   size=12, color=LIGHT_GRAY)
    rule(s, 4.28, color=ELECTRIC, l=l+0.3, w=3.25, h=0.03)
    for j, feat in enumerate(feats):
        box(s, l+0.15, 4.42 + j * 0.46, 3.65, 0.42, f"✓  {feat}", size=12, color=WHITE)

box(s, 0.4, 7.05, 12.5, 0.38,
    "Implementation fee: $2,500  (waived with 12-month commitment)  ·  Overage: $5/invoice  ·  60-day satisfaction guarantee",
    size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — NEXT STEPS
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s); left_stripe(s)

box(s, 0.55, 0.9,  12.0, 0.6,  "NEXT STEPS", size=11, color=ELECTRIC, bold=True)
box(s, 0.55, 1.4,  11.5, 1.0,
    "From First Conversation to Fully Automated — in Under 30 Days",
    size=30, bold=True, color=WHITE)
rule(s, 2.62)

steps_next = [
    ("Today",             "Scoping conversation — map your current cash application process, volume, and ERP setup"),
    ("Within 48 Hours",   "Custom proposal with recommended plan tier, integration plan, and projected ROI"),
    ("Week 1–2",          "Technical onboarding — bank feed connection, ERP integration, matching rules configured"),
    ("Day 14",            "Go live — your team receives their first automated cash application summary"),
    ("Day 60",            "Satisfaction checkpoint — 80% reduction in manual time or you pay nothing"),
]
for i, (when, what) in enumerate(steps_next):
    t = 2.88 + i * 0.85
    dot(s, 0.55, t+0.08)
    box(s, 0.92, t,      2.5,  0.42, when, size=14, bold=True, color=ELECTRIC)
    box(s, 0.92, t+0.42, 11.5, 0.38, what, size=13, color=LIGHT_GRAY)

card(s, 7.55, 2.78, 5.4, 4.05)
box(s, 7.75, 2.95, 5.0, 0.5,  "JONATHAN RODRIGUEZ",            size=16, bold=True, color=WHITE)
box(s, 7.75, 3.45, 5.0, 0.45, "Founder & CEO, LunarLogic",     size=13, color=LIGHT_GRAY)
rule(s, 4.08, color=ELECTRIC, l=7.75, w=4.8, h=0.03)
box(s, 7.75, 4.22, 5.0, 0.45, "✉  jrodriguez@lunarlogic.ai",   size=14, color=WHITE)
rule(s, 4.88, color=DARK_CARD, l=7.75, w=4.8, h=0.03)
box(s, 7.75, 5.0,  5.0, 0.45, "60-Day Satisfaction Guarantee", size=14, bold=True, color=GREEN)
box(s, 7.75, 5.45, 5.0, 0.88,
    "If we don't reduce your cash application\ntime by 80%, you pay nothing.",
    size=13, color=LIGHT_GRAY)

out = "/home/user/lunarlogic-dashboard/LunarLogic_CashApplication_PeterSukits_v2.pptx"
prs.save(out)
print(f"Saved → {out}")
